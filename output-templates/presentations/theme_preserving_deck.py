#!/usr/bin/env python3
"""
Install: python3 -m pip install python-pptx

Inject a normalized slide plan into an existing PPTX/POTX without changing its
theme, masters, layouts, colors, fonts, backgrounds, or brand assets.

The script only adds slides through template-owned slide layouts and only writes
into placeholders already present on those layouts. It never creates a text box,
shape, theme, master, layout, background, color, font, or image.

Usage:
    python3 theme_preserving_deck.py \
      --template branded-template.pptx \
      --plan slide-plan.json \
      --output injected-deck.pptx \
      --audit-output injection-audit.json
"""
from __future__ import annotations

import argparse
import json
import sys
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER


TITLE_TYPES = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}
BODY_TYPES = {
    PP_PLACEHOLDER.BODY,
    PP_PLACEHOLDER.OBJECT,
    PP_PLACEHOLDER.VERTICAL_BODY,
    PP_PLACEHOLDER.VERTICAL_OBJECT,
}
SUBTITLE_TYPES = {PP_PLACEHOLDER.SUBTITLE}


@dataclass(frozen=True)
class SlideSpec:
    role: str
    title: str
    subtitle: str
    bullets: tuple[str, ...]


def fail(message: str) -> "None":
    print(json.dumps({"error": message}))
    raise SystemExit(1)


def load_plan(path: Path) -> list[SlideSpec]:
    try:
        payload = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as exc:
        fail(f"could not read slide plan: {exc}")
    raw_slides = payload.get("slides") if isinstance(payload, dict) else None
    if not isinstance(raw_slides, list) or not raw_slides:
        fail("slide plan must contain a non-empty 'slides' array")
    slides: list[SlideSpec] = []
    allowed_roles = {"title", "section", "content", "summary"}
    for index, raw in enumerate(raw_slides, start=1):
        if not isinstance(raw, dict):
            fail(f"slide {index} must be an object")
        role = str(raw.get("role", "content")).strip().lower()
        if role not in allowed_roles:
            fail(f"slide {index} role must be one of {sorted(allowed_roles)}")
        title = str(raw.get("title", "")).strip()
        subtitle = str(raw.get("subtitle", "")).strip()
        raw_bullets = raw.get("bullets", [])
        if not isinstance(raw_bullets, list):
            fail(f"slide {index} bullets must be an array")
        bullets = tuple(str(item).strip() for item in raw_bullets if str(item).strip())
        if not title:
            fail(f"slide {index} needs a title")
        slides.append(SlideSpec(role, title, subtitle, bullets))
    return slides


def placeholder_types(layout) -> list[int]:
    return [shape.placeholder_format.type for shape in layout.placeholders]


def score_layout(layout, spec: SlideSpec) -> tuple[int, str]:
    name = (layout.name or "").lower()
    types = placeholder_types(layout)
    has_title = any(kind in TITLE_TYPES for kind in types)
    has_body = any(kind in BODY_TYPES for kind in types)
    has_subtitle = any(kind in SUBTITLE_TYPES for kind in types)
    score = 0
    reasons: list[str] = []

    if has_title:
        score += 40
        reasons.append("title-placeholder")
    else:
        score -= 100

    if spec.role == "title":
        if "title slide" in name or name == "title":
            score += 60
            reasons.append("title-layout-name")
        if has_subtitle:
            score += 30
            reasons.append("subtitle-placeholder")
        if has_body:
            score -= 10
    elif spec.role == "section":
        if "section" in name:
            score += 60
            reasons.append("section-layout-name")
        if has_subtitle or has_body:
            score += 20
    else:
        if "title and content" in name or "content" in name:
            score += 50
            reasons.append("content-layout-name")
        if has_body:
            score += 50
            reasons.append("body-placeholder")
        elif spec.bullets:
            score -= 100

    if "blank" in name:
        score -= 80
    return score, ", ".join(reasons) or "nearest available layout"


def choose_layout(prs: Presentation, spec: SlideSpec):
    ranked = []
    for index, layout in enumerate(prs.slide_layouts):
        score, reason = score_layout(layout, spec)
        ranked.append((score, -index, layout, reason))
    ranked.sort(key=lambda item: (item[0], item[1]), reverse=True)
    score, _, layout, reason = ranked[0]
    if score < 0:
        needed = "title and body" if spec.bullets else "title"
        fail(f"template has no usable {needed} placeholder layout for role '{spec.role}'")
    return layout, score, reason


def placeholders_by_type(slide, accepted: set[int]):
    return [
        shape for shape in slide.placeholders
        if shape.placeholder_format.type in accepted and getattr(shape, "has_text_frame", False)
    ]


def set_plain_text(placeholder, text: str) -> None:
    frame = placeholder.text_frame
    frame.clear()
    frame.paragraphs[0].text = text


def set_bullets(placeholder, bullets: tuple[str, ...]) -> None:
    frame = placeholder.text_frame
    frame.clear()
    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0


def populate_slide(slide, spec: SlideSpec) -> dict:
    title_slots = placeholders_by_type(slide, TITLE_TYPES)
    if not title_slots:
        fail(f"selected layout for '{spec.title}' has no title placeholder")
    set_plain_text(title_slots[0], spec.title)
    used = [{"idx": title_slots[0].placeholder_format.idx, "purpose": "title"}]
    fallbacks: list[str] = []

    subtitle_slots = placeholders_by_type(slide, SUBTITLE_TYPES)
    body_slots = placeholders_by_type(slide, BODY_TYPES)

    if spec.subtitle:
        if subtitle_slots:
            set_plain_text(subtitle_slots[0], spec.subtitle)
            used.append({"idx": subtitle_slots[0].placeholder_format.idx, "purpose": "subtitle"})
        elif body_slots and not spec.bullets:
            set_plain_text(body_slots[0], spec.subtitle)
            used.append({"idx": body_slots[0].placeholder_format.idx, "purpose": "subtitle-fallback"})
            fallbacks.append("subtitle placed in the template's body placeholder")
        else:
            fallbacks.append("subtitle omitted because the selected layout has no unused text placeholder")

    if spec.bullets:
        available = list(body_slots)
        if spec.subtitle and not subtitle_slots and available:
            available = available[1:]
        if not available:
            fail(f"selected layout for '{spec.title}' has no body placeholder for supplied bullets")
        set_bullets(available[0], spec.bullets)
        used.append({"idx": available[0].placeholder_format.idx, "purpose": "bullets"})

    return {"placeholders_written": used, "fallbacks": fallbacks}


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--template", required=True)
    parser.add_argument("--plan", required=True)
    parser.add_argument("--output", required=True)
    parser.add_argument("--audit-output", required=True)
    args = parser.parse_args()

    template = Path(args.template).resolve()
    plan_path = Path(args.plan).resolve()
    output = Path(args.output).resolve()
    audit_path = Path(args.audit_output).resolve()
    if template.suffix.lower() not in {".pptx", ".potx"}:
        fail("template must be a .pptx or .potx file")
    if not template.exists():
        fail(f"template not found: {template}")

    specs = load_plan(plan_path)
    try:
        prs = Presentation(str(template))
    except Exception as exc:
        fail(f"could not open template: {exc}")

    starting_slides = len(prs.slides)
    layout_inventory = [
        {
            "index": index,
            "name": layout.name,
            "placeholder_types": [str(kind) for kind in placeholder_types(layout)],
        }
        for index, layout in enumerate(prs.slide_layouts)
    ]
    mappings: list[dict] = []
    for index, spec in enumerate(specs, start=1):
        layout, layout_score, reason = choose_layout(prs, spec)
        slide = prs.slides.add_slide(layout)
        populated = populate_slide(slide, spec)
        mappings.append({
            "plan_index": index,
            "role": spec.role,
            "title": spec.title,
            "layout_name": layout.name,
            "layout_score": layout_score,
            "selection_basis": reason,
            **populated,
        })

    output.parent.mkdir(parents=True, exist_ok=True)
    audit_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    audit = {
        "template": template.name,
        "plan": plan_path.name,
        "output": output.name,
        "starting_slides": starting_slides,
        "slides_added": len(specs),
        "ending_slides": len(prs.slides),
        "masters_before_save": len(prs.slide_masters),
        "layouts_before_save": len(prs.slide_layouts),
        "layout_inventory": layout_inventory,
        "mappings": mappings,
        "policy": {
            "slides_added_only_with_template_layouts": True,
            "content_written_only_to_placeholders": True,
            "new_shapes_created": False,
            "theme_or_master_edited": False,
        },
    }
    audit_path.write_text(json.dumps(audit, indent=2) + "\n", encoding="utf-8")
    print(json.dumps({
        "output": output.name,
        "slides_added": len(specs),
        "starting_slides": starting_slides,
        "ending_slides": len(prs.slides),
        "audit": audit_path.name,
    }, indent=2))
    return 0


if __name__ == "__main__":
    sys.exit(main())
