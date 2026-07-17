#!/usr/bin/env python3
"""Create the fictional Northstar Studio PPTX fixture used for preservation QA."""
from __future__ import annotations

import argparse
import copy
import tempfile
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from pptx import Presentation


NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
    "ct": "http://schemas.openxmlformats.org/package/2006/content-types",
}
for prefix in ("a", "p", "r"):
    ET.register_namespace(prefix, NS[prefix])
ET.register_namespace("", NS["rel"])


COLORS = {
    "dk1": "10243E",
    "lt1": "F7F4EE",
    "dk2": "203A5F",
    "lt2": "E8E1D5",
    "accent1": "D6A84B",
    "accent2": "4E7C86",
    "accent3": "8E6A5A",
    "accent4": "6F7D9B",
    "accent5": "7E9154",
    "accent6": "C26C4A",
    "hlink": "2B6CB0",
    "folHlink": "7C3AED",
}


def q(prefix: str, local: str) -> str:
    return f"{{{NS[prefix]}}}{local}"


def patch_theme(data: bytes) -> bytes:
    root = ET.fromstring(data)
    root.set("name", "Northstar Studio Theme")
    scheme = root.find(".//a:clrScheme", NS)
    if scheme is None:
        raise RuntimeError("template theme has no color scheme")
    scheme.set("name", "Northstar Studio")
    for name, value in COLORS.items():
        holder = scheme.find(f"a:{name}", NS)
        if holder is None or len(holder) == 0:
            raise RuntimeError(f"theme is missing color slot {name}")
        color = holder[0]
        color.tag = q("a", "srgbClr")
        color.attrib.clear()
        color.set("val", value)
    font_scheme = root.find(".//a:fontScheme", NS)
    if font_scheme is None:
        raise RuntimeError("template theme has no font scheme")
    font_scheme.set("name", "Northstar Studio")
    major = font_scheme.find("a:majorFont/a:latin", NS)
    minor = font_scheme.find("a:minorFont/a:latin", NS)
    if major is None or minor is None:
        raise RuntimeError("template theme has no Latin font declarations")
    major.set("typeface", "Georgia")
    minor.set("typeface", "Aptos")
    return ET.tostring(root, encoding="utf-8", xml_declaration=True)


def patch_master(data: bytes, image_rel_id: str) -> bytes:
    root = ET.fromstring(data)
    # The fixture owns these layout choices. Lower the title frame and use a
    # restrained template-defined type scale so two-line takeaway titles remain
    # inside the slide canvas without the injector restyling any placeholder.
    for shape_node in root.findall(".//p:sp", NS):
        placeholder = shape_node.find("p:nvSpPr/p:nvPr/p:ph", NS)
        if placeholder is not None and placeholder.get("type") == "title":
            offset = shape_node.find("p:spPr/a:xfrm/a:off", NS)
            extent = shape_node.find("p:spPr/a:xfrm/a:ext", NS)
            if offset is not None:
                offset.set("y", "457200")
            if extent is not None:
                extent.set("cy", "1051560")
    title_size = root.find("p:txStyles/p:titleStyle/a:lvl1pPr/a:defRPr", NS)
    if title_size is not None:
        title_size.set("sz", "3600")
    body_sizes = ["2400", "2200", "2000", "1800", "1800", "1800", "1800", "1800", "1800"]
    for level, size in enumerate(body_sizes, start=1):
        node = root.find(f"p:txStyles/p:bodyStyle/a:lvl{level}pPr/a:defRPr", NS)
        if node is not None:
            node.set("sz", size)
    tree = root.find(".//p:spTree", NS)
    if tree is None:
        raise RuntimeError("slide master has no shape tree")
    pic = ET.Element(q("p", "pic"))
    nv = ET.SubElement(pic, q("p", "nvPicPr"))
    ET.SubElement(nv, q("p", "cNvPr"), {"id": "900", "name": "Northstar Studio Logo"})
    cnv = ET.SubElement(nv, q("p", "cNvPicPr"))
    ET.SubElement(cnv, q("a", "picLocks"), {"noChangeAspect": "1"})
    ET.SubElement(nv, q("p", "nvPr"))
    fill = ET.SubElement(pic, q("p", "blipFill"))
    ET.SubElement(fill, q("a", "blip"), {q("r", "embed"): image_rel_id})
    stretch = ET.SubElement(fill, q("a", "stretch"))
    ET.SubElement(stretch, q("a", "fillRect"))
    shape = ET.SubElement(pic, q("p", "spPr"))
    xfrm = ET.SubElement(shape, q("a", "xfrm"))
    ET.SubElement(xfrm, q("a", "off"), {"x": "8275320", "y": "228600"})
    ET.SubElement(xfrm, q("a", "ext"), {"cx": "457200", "cy": "457200"})
    geom = ET.SubElement(shape, q("a", "prstGeom"), {"prst": "rect"})
    ET.SubElement(geom, q("a", "avLst"))
    tree.append(pic)
    return ET.tostring(root, encoding="utf-8", xml_declaration=True)


def patch_master_rels(data: bytes, image_rel_id: str) -> bytes:
    ET.register_namespace("", NS["rel"])
    root = ET.fromstring(data)
    ET.SubElement(root, f"{{{NS['rel']}}}Relationship", {
        "Id": image_rel_id,
        "Type": "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image",
        "Target": "../media/northstar-logo.png",
    })
    return ET.tostring(root, encoding="utf-8", xml_declaration=True)


def patch_content_types(data: bytes) -> bytes:
    ET.register_namespace("", NS["ct"])
    root = ET.fromstring(data)
    has_png = any(node.get("Extension", "").lower() == "png" for node in root)
    if not has_png:
        ET.SubElement(root, f"{{{NS['ct']}}}Default", {
            "Extension": "png", "ContentType": "image/png",
        })
    rendered = ET.tostring(root, encoding="utf-8", xml_declaration=True)
    ET.register_namespace("", NS["rel"])
    return rendered


def create_base(path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Northstar Studio"
    subtitle = next(
        shape for shape in slide.placeholders
        if shape.placeholder_format.type == 4 and shape.has_text_frame
    )
    subtitle.text = "Presentation template | Fictional brand-preservation fixture"
    prs.save(path)


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--logo", required=True)
    parser.add_argument("--output", required=True)
    args = parser.parse_args()
    logo = Path(args.logo).resolve()
    output = Path(args.output).resolve()
    if not logo.exists():
        raise SystemExit(f"logo not found: {logo}")
    output.parent.mkdir(parents=True, exist_ok=True)

    with tempfile.TemporaryDirectory(prefix="northstar-template-") as tmp:
        base = Path(tmp) / "base.pptx"
        create_base(base)
        image_rel_id = "rId900"
        with zipfile.ZipFile(base, "r") as source, zipfile.ZipFile(
            output, "w", compression=zipfile.ZIP_DEFLATED
        ) as target:
            for item in source.infolist():
                data = source.read(item.filename)
                if item.filename == "ppt/theme/theme1.xml":
                    data = patch_theme(data)
                elif item.filename == "ppt/slideMasters/slideMaster1.xml":
                    data = patch_master(data, image_rel_id)
                elif item.filename == "ppt/slideMasters/_rels/slideMaster1.xml.rels":
                    data = patch_master_rels(data, image_rel_id)
                elif item.filename == "[Content_Types].xml":
                    data = patch_content_types(data)
                target.writestr(item, data)
            target.writestr("ppt/media/northstar-logo.png", logo.read_bytes())
    print(f"created {output}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
