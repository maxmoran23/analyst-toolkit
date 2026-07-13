#!/usr/bin/env python3
"""Prove that a generated deck retained its template-owned brand package parts."""
from __future__ import annotations

import argparse
import hashlib
import json
import posixpath
import sys
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from pptx import Presentation


A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"


def sha(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def canonical_node(node) -> tuple:
    attrs = tuple(sorted(node.attrib.items()))
    text = (node.text or "").strip()
    children = tuple(canonical_node(child) for child in list(node))
    return node.tag, attrs, text, children


def canonical_hash(data: bytes) -> str:
    value = repr(canonical_node(ET.fromstring(data))).encode("utf-8")
    return sha(value)


def local(tag: str) -> str:
    return tag.rsplit("}", 1)[-1]


def theme_contract(data: bytes) -> dict:
    root = ET.fromstring(data)
    colors = {}
    fonts = {}
    for node in root.iter():
        name = local(node.tag)
        parent_slots = {
            "dk1", "lt1", "dk2", "lt2", "accent1", "accent2", "accent3",
            "accent4", "accent5", "accent6", "hlink", "folHlink",
        }
        if name in parent_slots and len(node):
            color = node[0]
            colors[name] = color.get("val") or color.get("lastClr")
    major = root.find(f".//{{{A_NS}}}majorFont/{{{A_NS}}}latin")
    minor = root.find(f".//{{{A_NS}}}minorFont/{{{A_NS}}}latin")
    fonts["major_latin"] = major.get("typeface") if major is not None else None
    fonts["minor_latin"] = minor.get("typeface") if minor is not None else None
    return {"colors": colors, "fonts": fonts, "raw_sha256": sha(data)}


def part_map(archive: zipfile.ZipFile, prefix: str) -> dict[str, dict]:
    return {
        name: {"sha256": sha(archive.read(name)), "canonical_sha256": canonical_hash(archive.read(name))}
        for name in sorted(archive.namelist())
        if name.startswith(prefix) and name.endswith((".xml", ".rels"))
    }


def master_images(archive: zipfile.ZipFile) -> list[dict]:
    found = []
    for rel_name in sorted(name for name in archive.namelist() if name.startswith(
        "ppt/slideMasters/_rels/") and name.endswith(".rels")):
        root = ET.fromstring(archive.read(rel_name))
        master_dir = "ppt/slideMasters"
        for rel in root:
            if rel.get("Type", "").endswith("/image"):
                target = posixpath.normpath(posixpath.join(master_dir, rel.get("Target", "")))
                payload = archive.read(target)
                found.append({
                    "relationship_part": rel_name,
                    "relationship_id": rel.get("Id"),
                    "target": target,
                    "sha256": sha(payload),
                })
    return found


def package_contract(path: Path) -> dict:
    with zipfile.ZipFile(path, "r") as archive:
        themes = sorted(name for name in archive.namelist() if name.startswith("ppt/theme/") and name.endswith(".xml"))
        return {
            "theme_parts": {name: theme_contract(archive.read(name)) for name in themes},
            "master_parts": part_map(archive, "ppt/slideMasters/"),
            "layout_parts": part_map(archive, "ppt/slideLayouts/"),
            "master_images": master_images(archive),
        }


def added_slide_contract(template: Path, output: Path, plan: Path) -> dict:
    before = Presentation(str(template))
    after = Presentation(str(output))
    payload = json.loads(plan.read_text(encoding="utf-8"))
    expected_strings = []
    for slide in payload["slides"]:
        expected_strings.extend([slide.get("title", ""), slide.get("subtitle", "")])
        expected_strings.extend(slide.get("bullets", []))
    expected_strings = [value for value in expected_strings if value]
    actual_text = []
    non_placeholder = []
    for slide_index in range(len(before.slides), len(after.slides)):
        slide = after.slides[slide_index]
        for shape in slide.shapes:
            if not shape.is_placeholder:
                non_placeholder.append({"slide": slide_index + 1, "shape": shape.name})
            if getattr(shape, "has_text_frame", False):
                actual_text.append(shape.text)
    joined = "\n".join(actual_text)
    missing = [value for value in expected_strings if value not in joined]
    return {
        "template_slide_count": len(before.slides),
        "output_slide_count": len(after.slides),
        "expected_added_slides": len(payload["slides"]),
        "actual_added_slides": len(after.slides) - len(before.slides),
        "non_placeholder_shapes_on_added_slides": non_placeholder,
        "missing_plan_text": missing,
    }


def compare(template: Path, output: Path, plan: Path) -> dict:
    before = package_contract(template)
    after = package_contract(output)
    slides = added_slide_contract(template, output, plan)
    def semantic_parts_equal(left: dict, right: dict) -> bool:
        return set(left) == set(right) and all(
            left[name]["canonical_sha256"] == right[name]["canonical_sha256"]
            for name in left
        )

    checks = {
        "theme_part_names_unchanged": set(before["theme_parts"]) == set(after["theme_parts"]),
        "theme_xml_byte_identical": before["theme_parts"] == after["theme_parts"],
        "theme_colors_and_fonts_unchanged": all(
            before["theme_parts"][name]["colors"] == after["theme_parts"].get(name, {}).get("colors")
            and before["theme_parts"][name]["fonts"] == after["theme_parts"].get(name, {}).get("fonts")
            for name in before["theme_parts"]
        ),
        "master_inventory_and_xml_unchanged": semantic_parts_equal(before["master_parts"], after["master_parts"]),
        "layout_inventory_and_xml_unchanged": semantic_parts_equal(before["layout_parts"], after["layout_parts"]),
        "master_logo_relationship_and_bytes_unchanged": before["master_images"] == after["master_images"] and bool(before["master_images"]),
        "expected_slide_count_added": slides["actual_added_slides"] == slides["expected_added_slides"],
        "all_added_slide_shapes_are_placeholders": not slides["non_placeholder_shapes_on_added_slides"],
        "all_planned_text_is_present": not slides["missing_plan_text"],
    }
    return {"checks": checks, "template": before, "output": after, "slides": slides, "passed": all(checks.values())}


def render_markdown(result: dict) -> str:
    lines = [
        "# Theme-Preservation Verification",
        "",
        "> **In plain terms:** The output deck was checked at the PowerPoint package level. "
        "The template's theme, master, layouts, fonts, colors, and master logo must be unchanged; only added slide content may differ.",
        "",
        "| Check | Result |",
        "| --- | --- |",
    ]
    for name, passed in result["checks"].items():
        lines.append(f"| `{name}` | {'PASS' if passed else 'FAIL'} |")
    lines.extend([
        "",
        "## Brand contract",
        "",
        f"- Theme colors: `{json.dumps(next(iter(result['template']['theme_parts'].values()))['colors'], sort_keys=True)}`",
        f"- Theme fonts: `{json.dumps(next(iter(result['template']['theme_parts'].values()))['fonts'], sort_keys=True)}`",
        f"- Master logo SHA-256: `{result['template']['master_images'][0]['sha256'] if result['template']['master_images'] else 'MISSING'}`",
        f"- Slides: {result['slides']['template_slide_count']} template + {result['slides']['actual_added_slides']} added = {result['slides']['output_slide_count']} output",
        "",
        "## Conclusion",
        "",
        "The theme-preservation law passed." if result["passed"] else "The theme-preservation law failed; do not deliver this deck.",
        "",
        f"**Confidence: {'HIGH' if result['passed'] else 'LOW'} — package-level comparisons and placeholder inspection {'all passed' if result['passed'] else 'identified a preservation breach'}.**",
        "",
    ])
    return "\n".join(lines)


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--template", required=True)
    parser.add_argument("--output", required=True)
    parser.add_argument("--plan", required=True)
    parser.add_argument("--json", required=True)
    parser.add_argument("--report", required=True)
    args = parser.parse_args()
    template = Path(args.template).resolve()
    output = Path(args.output).resolve()
    plan = Path(args.plan).resolve()
    result = compare(template, output, plan)
    json_path = Path(args.json).resolve()
    report_path = Path(args.report).resolve()
    json_path.parent.mkdir(parents=True, exist_ok=True)
    report_path.parent.mkdir(parents=True, exist_ok=True)
    json_path.write_text(json.dumps(result, indent=2) + "\n", encoding="utf-8")
    report_path.write_text(render_markdown(result), encoding="utf-8")
    print(json.dumps({"passed": result["passed"], "checks": result["checks"]}, indent=2))
    return 0 if result["passed"] else 1


if __name__ == "__main__":
    sys.exit(main())
